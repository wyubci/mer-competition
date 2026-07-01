# -*- coding: utf-8 -*-
"""Add explicit slides explaining video/time prior vs direct physiology routes."""

from __future__ import annotations

import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPTX = Path(os.environ.get("MERPS_PPTX", str(ROOT / "MER_PS_strategy_report.pptx")))
OUTLINE = Path(os.environ.get("MERPS_OUTLINE", str(ROOT / "MER_PS_strategy_report_outline.md")))
RESULT_DIR = ROOT / "experiments" / "results"

FONT = "Microsoft YaHei"
BG = RGBColor(248, 250, 252)
INK = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
PURPLE = RGBColor(124, 58, 237)
LINE = RGBColor(203, 213, 225)
WHITE = RGBColor(255, 255, 255)


def set_font(run, size=None, bold=False, color=INK):
    run.font.name = FONT
    run.font.bold = bold
    run.font.color.rgb = color
    if size is not None:
        run.font.size = Pt(size)


def add_bg(slide):
    rect = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG
    rect.line.fill.background()
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.1), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = title
    set_font(r, 24, True, INK)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.96), Inches(12.0), Inches(0.32))
        stf = sub.text_frame
        stf.clear()
        r2 = stf.paragraphs[0].add_run()
        r2.text = subtitle
        set_font(r2, 10.5, False, MUTED)


def add_footer(slide, idx):
    box = slide.shapes.add_textbox(Inches(11.8), Inches(7.06), Inches(0.9), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{idx:02d}"
    set_font(r, 8, False, MUTED)


def add_text(slide, x, y, w, h, text, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, bold, color)


def add_box(slide, x, y, w, h, title, body, color):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.1)
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    set_font(r, 12.4, True, color)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = body
    set_font(r2, 9.4, False, MUTED)


def add_table(slide, x, y, w, h, headers, rows, widths=None, font_size=8.8):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = ""
        r = p.add_run()
        r.text = header
        set_font(r, font_size, True, WHITE)
    for r_i, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            cell = table.cell(r_i, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_i % 2 else RGBColor(241, 245, 249)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            p.text = ""
            r = p.add_run()
            r.text = str(value)
            set_font(r, font_size, False, INK)


def new_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle)
    return slide


def load_results():
    best = json.loads((RESULT_DIR / "iteration_221_228_bcrf_module_seed2026.json").read_text(encoding="utf-8"))[
        "aggregate_results"
    ][0]
    no_prior = json.loads((RESULT_DIR / "iteration_321_335_no_video_prior_signal.json").read_text(encoding="utf-8"))[
        "aggregate_results"
    ][0]
    return best, no_prior


def add_prior_slide(prs):
    slide = new_slide(
        prs,
        "为什么使用视频/时间先验",
        "它不是测试标签泄漏，而是同一情绪刺激在跨主体任务中的稳定公共轨迹",
    )
    add_box(
        slide,
        0.75,
        1.45,
        3.25,
        1.1,
        "赛题结构",
        "每个主体观看同一组 15 段视频；sample_id 包含 Vxx 与 Txxx。",
        BLUE,
    )
    add_box(
        slide,
        4.35,
        1.45,
        3.25,
        1.1,
        "情绪来源",
        "动态 valence/arousal 很大一部分由视频刺激和播放时间驱动。",
        TEAL,
    )
    add_box(
        slide,
        7.95,
        1.45,
        3.65,
        1.1,
        "合法边界",
        "只用训练主体标签估计 p0(v,t)，不读取测试标签。",
        GREEN,
    )
    add_text(
        slide,
        0.9,
        3.05,
        11.55,
        0.45,
        "核心公式：p0(v,t) = E_train[y | video=v, time=t]",
        18,
        True,
        INK,
        PP_ALIGN.CENTER,
    )
    add_table(
        slide,
        0.8,
        3.9,
        11.85,
        1.75,
        ["原因", "解释", "对模型的影响"],
        [
            ["刺激一致", "同一视频在同一秒通常诱发相似情绪趋势", "形成强低维基线"],
            ["跨主体噪声大", "EEG/fNIRS 个体差异和噪声都很强", "先验可降低预测方差"],
            ["比赛目标", "输入 sample_id 本身给出视频和时间位置", "适合冲榜，但需说明先验属性"],
        ],
        widths=[1.65, 5.15, 5.05],
        font_size=8.8,
    )
    add_text(
        slide,
        0.95,
        6.0,
        11.35,
        0.65,
        "因此，视频/时间先验是“刺激轨迹建模”：它解决公共情绪轨迹；生理信号解决个体偏移和局部残差。",
        12,
        True,
        AMBER,
    )


def add_difference_slide(prs):
    best, no_prior = load_results()
    slide = new_slide(
        prs,
        "两条提交路线的区别",
        "A 追求比赛分数，B 验证纯生理证据；两者回答的是不同科学问题",
    )
    add_table(
        slide,
        0.65,
        1.42,
        12.05,
        3.75,
        ["维度", "路线 A：视频/时间先导", "路线 B：直接生理信号"],
        [
            ["输入", "sample_id + EEG/fNIRS；可使用 video/time prior", "只使用 EEG/fNIRS 特征"],
            ["预测对象", "公共刺激轨迹 + 个体残差", "生理信号到情绪标签的直接映射"],
            ["当前最好", f"{best['overall_mae']:.4f} MAE", f"{no_prior['overall_mae']:.4f} MAE"],
            ["优势", "分数强、稳定、符合当前 leaderboard 目标", "更干净，能回答生理信号是否真的可解码"],
            ["风险", "可能被质疑依赖视频刺激模式", "分数弱，容易输给中心/官方基线附近"],
            ["论文角色", "主方法性能上限与工程提交", "关键消融：去先验生理证据"],
        ],
        widths=[1.75, 5.15, 5.15],
        font_size=8.6,
    )
    add_text(
        slide,
        0.9,
        5.65,
        11.45,
        0.75,
        "一句话区别：路线 A 预测“某视频此刻大家大概率是什么情绪，再用生理信号微调”；路线 B 预测“只看脑电/近红外能不能解码情绪”。",
        12.2,
        True,
        INK,
    )


def add_synergy_slide(prs):
    best, no_prior = load_results()
    slide = new_slide(
        prs,
        "二者是否合力：是，但要分层融合",
        "不是简单平均两个模型，而是让先验、校准、生理残差各自承担不同误差项",
    )
    add_text(
        slide,
        0.8,
        1.45,
        11.9,
        0.55,
        "ŷ = clip( p0(v,t) + Δ_cal(SCRF-BCRF) + λ · g_ccmi · Δ_sig(EEG,fNIRS), 1, 255 )",
        17,
        True,
        GREEN,
        PP_ALIGN.CENTER,
    )
    add_table(
        slide,
        0.75,
        2.35,
        11.9,
        2.25,
        ["项", "来自哪条路线", "解决什么问题", "为什么不能互相替代"],
        [
            ["p0(v,t)", "视频/时间先导", "公共视频刺激的平均情绪轨迹", "生理信号太噪，直接学公共轨迹不稳"],
            ["Δ_cal", "SCRF-BCRF", "先验的系统性偏差", "直接信号模型没有稳定估计可信度"],
            ["Δ_sig", "直接生理路线", "个体生理状态与局部偏移", "视频先验看不到个体生理差异"],
            ["λ·g_ccmi", "合力门控", "只在 EEG 与 fNIRS 同向时放大生理残差", "防止单模态噪声破坏强基线"],
        ],
        widths=[1.45, 2.25, 3.95, 4.25],
        font_size=8.5,
    )
    add_box(
        slide,
        0.95,
        5.25,
        3.55,
        0.95,
        "主提交",
        f"路线 A：{best['overall_mae']:.4f} MAE\n用于冲榜",
        BLUE,
    )
    add_box(
        slide,
        4.85,
        5.25,
        3.55,
        0.95,
        "消融/备选提交",
        f"路线 B：{no_prior['overall_mae']:.4f} MAE\n用于证明生理证据",
        TEAL,
    )
    add_box(
        slide,
        8.75,
        5.25,
        3.55,
        0.95,
        "最终合力方向",
        "p0 + Δ_cal + 小权重 CCMI\n仍需 per-sample cache 验证",
        PURPLE,
    )


def update_outline(slide_count):
    marker = "Video-prior vs direct-physiology logic slides:"
    old = OUTLINE.read_text(encoding="utf-8") if OUTLINE.exists() else ""
    if marker in old:
        old = old.split(marker, 1)[0].rstrip() + "\n"
    lines = [
        "",
        marker,
        f"{slide_count - 2}. 为什么使用视频/时间先验",
        f"{slide_count - 1}. 两条提交路线的区别",
        f"{slide_count}. 二者是否合力：分层融合公式",
    ]
    OUTLINE.write_text(old.rstrip() + "\n" + "\n".join(lines) + "\n", encoding="utf-8")


def main():
    prs = Presentation(PPTX)
    add_prior_slide(prs)
    add_difference_slide(prs)
    add_synergy_slide(prs)
    for idx, slide in enumerate(prs.slides, 1):
        if idx > len(prs.slides) - 3:
            add_footer(slide, idx)
    prs.save(PPTX)
    update_outline(len(prs.slides))
    print(f"updated {PPTX}")
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    main()
