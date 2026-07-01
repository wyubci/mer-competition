# -*- coding: utf-8 -*-
"""Append innovation-module diagrams to the MER-PS strategy deck.

Run after tools/make_merps_strategy_ppt.py. The script updates the existing
deck in place and also refreshes the outline with the added detail slides.
"""

from __future__ import annotations

import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
RESULT_DIR = ROOT / "experiments" / "results"
PPTX = Path(os.environ.get("MERPS_PPTX", str(ROOT / "MER_PS_strategy_report.pptx")))
OUTLINE = Path(os.environ.get("MERPS_OUTLINE", str(ROOT / "MER_PS_strategy_report_outline.md")))

FONT = "Microsoft YaHei"
BG = RGBColor(248, 250, 252)
INK = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
PURPLE = RGBColor(124, 58, 237)
LINE = RGBColor(203, 213, 225)
WHITE = RGBColor(255, 255, 255)


def set_font(run, size: float | None = None, bold: bool = False, color: RGBColor = INK) -> None:
    run.font.name = FONT
    run.font.bold = bold
    run.font.color.rgb = color
    if size is not None:
        run.font.size = Pt(size)


def add_bg(slide) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG
    rect.line.fill.background()
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.72))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, 25, True, INK)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.0), Inches(11.7), Inches(0.34))
        stf = sub.text_frame
        stf.clear()
        r2 = stf.paragraphs[0].add_run()
        r2.text = subtitle
        set_font(r2, 10.5, False, MUTED)


def add_footer(slide, idx: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.85), Inches(7.05), Inches(0.9), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{idx:02d}"
    set_font(r, 8, False, MUTED)


def add_text(slide, x: float, y: float, w: float, h: float, text: str, size: float = 12,
             bold: bool = False, color: RGBColor = INK, align=PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, bold, color)


def add_box(slide, x: float, y: float, w: float, h: float, title: str, body: str = "",
            color: RGBColor = BLUE, fill: RGBColor = WHITE, title_size: float = 11.5,
            body_size: float = 8.6):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = color
    shp.line.width = Pt(1.4)
    tf = shp.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    set_font(r, title_size, True, color)
    if body:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = body
        set_font(r2, body_size, False, MUTED)
    return shp


def connect(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = LINE, width: float = 1.4):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def add_table(slide, x: float, y: float, w: float, h: float, headers: list[str], rows: list[list[str]],
              widths: list[float] | None = None, font_size: float = 8.5) -> None:
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = ""
        r = p.add_run()
        r.text = header
        set_font(r, font_size, True, WHITE)
    for r_i, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            cell = table.cell(r_i, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_i % 2 else RGBColor(241, 245, 249)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            p.text = ""
            r = p.add_run()
            r.text = value
            set_font(r, font_size, False, INK)


def replace_exact_text(prs: Presentation) -> None:
    replacements = {
        "目标：用 EEG + fNIRS 预测 1 Hz 连续 valence / arousal 轨迹":
            "目标：用 EEG（脑电图）+ fNIRS（功能性近红外光谱）预测 1 Hz 连续 valence（效价）/ arousal（唤醒度）轨迹",
        "当前最佳：222_BCRF_onSCRF，Overall MAE = 28.6868。":
            "当前最佳：222_BCRF_onSCRF，其中 BCRF（贝叶斯可信残差场）叠加 SCRF（符号一致残差场），Overall MAE（总体平均绝对误差）= 28.6868。",
        "最佳 EEG-fNIRS 信号融合：CCMI PriorSlopeGate，Overall MAE = 28.7145。":
            "最佳 EEG-fNIRS 信号融合：CCMI（保守跨模态交集） PriorSlopeGate，Overall MAE = 28.7145。",
        "MER-PS: synchronized EEG + fNIRS continuous emotion decoding":
            "MER-PS: synchronized EEG（脑电图）+ fNIRS（功能性近红外光谱） continuous emotion decoding（连续情绪解码）",
        "test_1 ~ test_24，subject-disjoint CV":
            "test_1 ~ test_24，subject-disjoint CV（跨主体交叉验证）",
        "15 段情绪诱发视频，含 5 秒 baseline":
            "15 段情绪诱发视频，含 5 秒 baseline（基线）",
        "Official ASAC demo":
            "Official ASAC demo（官方 ASAC 演示基线）",
        "PatternPrior_098":
            "PatternPrior（视频-时间模式先验）_098",
    }
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)


def new_slide(prs: Presentation, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle)
    return slide


def slide_overall_module(prs: Presentation) -> None:
    s = new_slide(
        prs,
        "原创模块总览：CRF-Fusion（可信残差融合）",
        "把强先验、可信校准和 EEG-fNIRS 生理证据拆成三条可解释支路",
    )
    add_text(s, 0.7, 1.35, 3.2, 0.3, "信号支路", 12, True, TEAL)
    add_text(s, 0.7, 3.0, 3.2, 0.3, "先验支路", 12, True, BLUE)
    add_text(s, 0.7, 4.65, 3.2, 0.3, "可信校准支路", 12, True, PURPLE)

    add_box(s, 0.75, 1.75, 1.75, 0.72, "EEG / fNIRS", "原始生理信号", TEAL)
    add_box(s, 2.9, 1.75, 1.95, 0.72, "φ_s(signal)", "频带、血氧、滞后", TEAL)
    add_box(s, 5.25, 1.75, 2.05, 0.72, "Δ_sig", "信号残差候选", TEAL)
    connect(s, 2.5, 2.11, 2.9, 2.11, TEAL)
    connect(s, 4.85, 2.11, 5.25, 2.11, TEAL)

    add_box(s, 0.75, 3.4, 1.75, 0.72, "sample_id", "视频与秒级位置", BLUE)
    add_box(s, 2.9, 3.4, 1.95, 0.72, "p0(v,t)", "视频-时间先验", BLUE)
    add_box(s, 5.25, 3.4, 2.05, 0.72, "ŷ0", "基础预测", BLUE)
    connect(s, 2.5, 3.76, 2.9, 3.76, BLUE)
    connect(s, 4.85, 3.76, 5.25, 3.76, BLUE)

    add_box(s, 0.75, 5.05, 1.75, 0.72, "OOF residual", "r = y - ŷ0", PURPLE)
    add_box(s, 2.9, 5.05, 1.95, 0.72, "SCRF", "符号一致残差场", PURPLE)
    add_box(s, 5.25, 5.05, 2.05, 0.72, "BCRF", "贝叶斯可信权重", PURPLE)
    connect(s, 2.5, 5.41, 2.9, 5.41, PURPLE)
    connect(s, 4.85, 5.41, 5.25, 5.41, PURPLE)

    add_box(
        s,
        8.15,
        2.55,
        3.95,
        1.35,
        "融合方程",
        "ŷ = clip(p0 + Δ_cal + λ·g_ccmi·Δ_sig, 1, 255)",
        GREEN,
        title_size=13,
        body_size=10.2,
    )
    add_box(s, 8.75, 4.55, 2.75, 0.82, "输出", "valence（效价） / arousal（唤醒度）", GREEN)
    connect(s, 7.3, 2.11, 8.15, 3.05, TEAL)
    connect(s, 7.3, 3.76, 8.15, 3.18, BLUE)
    connect(s, 7.3, 5.41, 8.15, 3.45, PURPLE)
    connect(s, 10.1, 3.9, 10.1, 4.55, GREEN)
    add_text(
        s,
        0.9,
        6.25,
        11.55,
        0.45,
        "核心思想：先验负责解释可重复的视频-时间轨迹，SCRF/BCRF 负责修正系统性偏差，CCMI 只在 EEG 与 fNIRS 同向支持时加入小幅生理残差。",
        11.5,
        False,
        INK,
    )


def slide_scrf_bcrf(prs: Presentation) -> None:
    s = new_slide(
        prs,
        "SCRF-BCRF（符号一致贝叶斯可信残差场）结构图",
        "将残差学习改造成低维、可置信、可解释的统计校准",
    )
    add_box(s, 0.75, 1.55, 2.0, 0.78, "输入", "OOF（折外预测）残差\nr_i = y_i - p_i", BLUE)
    add_box(s, 3.1, 1.55, 2.35, 0.78, "状态分组", "video-time / time-value\nvalue-slope / video", TEAL)
    add_box(s, 5.8, 1.55, 2.3, 0.78, "稳健统计", "median, MAD, n, z\nsign_strength", AMBER)
    add_box(s, 8.45, 1.55, 2.35, 0.78, "可信权重", "w_g = shrink(n) · robust(MAD)\n· sigmoid(z-τ) · sign", PURPLE)
    add_box(s, 11.1, 1.55, 1.55, 0.78, "输出", "Δ_cal", GREEN)
    for x1, x2 in [(2.75, 3.1), (5.45, 5.8), (8.1, 8.45), (10.8, 11.1)]:
        connect(s, x1, 1.94, x2, 1.94)

    add_box(
        s,
        0.9,
        3.0,
        3.55,
        1.25,
        "SCRF（符号一致残差场）",
        "只保留多个低维视角方向一致的残差。\n如果某个 cell 内符号分裂，直接降低修正强度。",
        TEAL,
        title_size=12.5,
        body_size=9.2,
    )
    add_box(
        s,
        4.9,
        3.0,
        3.55,
        1.25,
        "BCRF（贝叶斯可信残差场）",
        "用样本数、MAD（中位绝对偏差）和 z-score（标准化置信度）估计 correction 是否可信。",
        PURPLE,
        title_size=12.5,
        body_size=9.2,
    )
    add_box(
        s,
        8.9,
        3.0,
        3.35,
        1.25,
        "Valence-only（只修效价）",
        "实验显示 arousal（唤醒度）残差不稳定，因此默认只对 valence 做小幅校准。",
        GREEN,
        title_size=12.5,
        body_size=9.2,
    )

    add_box(
        s,
        1.15,
        5.15,
        10.95,
        0.95,
        "最终形式",
        "Δ_cal(i) = Σ_g  w_g(i) · median(r | cell_g(i))，并经过符号一致门控与幅度裁剪",
        RED,
        title_size=12.5,
        body_size=10,
    )
    add_text(
        s,
        1.0,
        6.35,
        11.4,
        0.45,
        "可发表点：模块不是简单残差叠加，而是把“是否应该修正”显式建模，适合小样本、跨主体、高噪声的生理情绪预测。",
        11.3,
        False,
        INK,
    )


def slide_ccmi(prs: Presentation) -> None:
    s = new_slide(
        prs,
        "CCMI（保守跨模态交集）融合图",
        "用 EEG 与 fNIRS 的同向交集约束生理残差，避免单模态过度修正",
    )
    eeg = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(1.8), Inches(3.2), Inches(2.2))
    eeg.fill.solid()
    eeg.fill.fore_color.rgb = RGBColor(219, 234, 254)
    eeg.fill.transparency = 25
    eeg.line.color.rgb = BLUE
    fnirs = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.25), Inches(1.8), Inches(3.2), Inches(2.2))
    fnirs.fill.solid()
    fnirs.fill.fore_color.rgb = RGBColor(204, 251, 241)
    fnirs.fill.transparency = 25
    fnirs.line.color.rgb = TEAL
    add_text(s, 1.55, 2.35, 1.2, 0.35, "EEG expert", 12, True, BLUE, PP_ALIGN.CENTER)
    add_text(s, 4.75, 2.35, 1.2, 0.35, "fNIRS expert", 12, True, TEAL, PP_ALIGN.CENTER)
    add_text(s, 3.15, 2.55, 1.35, 0.35, "同向交集", 12, True, AMBER, PP_ALIGN.CENTER)
    add_text(s, 1.55, 2.85, 1.25, 0.35, "e_i", 14, True, BLUE, PP_ALIGN.CENTER)
    add_text(s, 4.85, 2.85, 1.25, 0.35, "f_i", 14, True, TEAL, PP_ALIGN.CENTER)

    add_box(s, 7.3, 1.55, 2.15, 0.75, "符号一致", "A = 1[sign(e)=sign(f)]", BLUE)
    add_box(s, 9.9, 1.55, 2.15, 0.75, "幅度交集", "m = min(|e|, |f|)", TEAL)
    add_box(s, 7.3, 2.85, 2.15, 0.75, "状态门控", "g_s = gate(slope(p0))", PURPLE)
    add_box(s, 9.9, 2.85, 2.15, 0.75, "方向", "d = sign(e+f)", AMBER)
    add_box(
        s,
        7.65,
        4.4,
        4.05,
        0.95,
        "CCMI residual",
        "Δ_sig = A · g_s · d · m",
        GREEN,
        title_size=13,
        body_size=11,
    )
    connect(s, 4.55, 2.9, 7.3, 1.92, AMBER)
    connect(s, 8.38, 2.3, 8.38, 2.85)
    connect(s, 10.98, 2.3, 10.98, 2.85)
    connect(s, 8.38, 3.6, 8.95, 4.4)
    connect(s, 10.98, 3.6, 10.35, 4.4)
    add_text(
        s,
        1.0,
        5.0,
        5.8,
        0.95,
        "设计依据：EEG 响应快但噪声高，fNIRS 响应慢但反映血氧变化。只有两者方向一致时，才认为生理证据足够可信。",
        11.2,
        False,
        INK,
    )
    add_text(
        s,
        7.25,
        5.75,
        4.85,
        0.62,
        "实验结论：最佳信号融合路径为 273_CCMI_PriorSlopeGate，Overall MAE = 28.7145。",
        11.2,
        True,
        GREEN,
    )


def slide_math(prs: Presentation, best: dict, no_prior_best: dict | None) -> None:
    s = new_slide(
        prs,
        "为什么这样拼接：数学逻辑",
        "把不同模块放在互补位置，而不是简单堆叠",
    )
    add_box(s, 0.8, 1.55, 3.6, 1.0, "标签分解", "y = p0(v,t) + r*", BLUE, title_size=13, body_size=12)
    add_box(s, 4.85, 1.55, 3.6, 1.0, "残差分解", "r* = r_bias + r_signal + ε", TEAL, title_size=13, body_size=12)
    add_box(s, 8.9, 1.55, 3.45, 1.0, "预测组合", "ŷ = p0 + Δ_cal + λΔ_sig", GREEN, title_size=13, body_size=12)
    connect(s, 4.4, 2.05, 4.85, 2.05)
    connect(s, 8.45, 2.05, 8.9, 2.05)

    rows = [
        ["p0(v,t)", "视频-时间先验", "解释视频刺激和时间轨迹的稳定均值"],
        ["Δ_cal", "SCRF-BCRF 校准", "修正先验在某些状态下的系统性偏差"],
        ["Δ_sig", "CCMI 生理残差", "只在 EEG 与 fNIRS 同向时加入小幅信号证据"],
        ["λ / gate", "缩放与门控", "控制自由度，减少跨主体过拟合"],
    ]
    add_table(s, 0.8, 3.0, 11.55, 2.1, ["符号", "模块", "作用"], rows, widths=[1.55, 2.55, 7.45], font_size=9.3)
    add_text(
        s,
        0.95,
        5.55,
        11.25,
        0.75,
        f"当前结果支持这个分解：全局最佳 {best['overall_mae']:.4f} 主要来自 p0 + Δ_cal；"
        + (
            f"no-prior 直接信号最佳 {no_prior_best['overall_mae']:.4f}，说明 Δ_sig 有信息但必须保守使用。"
            if no_prior_best
            else "下一步需要用 no-prior 实验继续验证 Δ_sig 的独立贡献。"
        ),
        11.3,
        False,
        INK,
    )


def slide_speaker_notes(prs: Presentation, best: dict, signal_best: dict, no_prior_best: dict | None) -> None:
    s = new_slide(
        prs,
        "汇报话术：贡献、证据与风险",
        "把比赛成绩和论文创新分开讲，避免被审稿人认为只是利用数据先验",
    )
    rows = [
        ["问题", "跨主体 EEG+fNIRS 连续情绪回归，小样本、高噪声、强视频时间结构"],
        ["核心贡献", "CRF-Fusion：可信残差校准 + 保守跨模态交集"],
        ["最好结果", f"222_BCRF_onSCRF: Overall MAE {best['overall_mae']:.4f}"],
        ["信号证据", f"CCMI PriorSlopeGate: Overall MAE {signal_best['overall_mae']:.4f}"],
        ["去先验结果", f"{no_prior_best['overall_mae']:.4f}" if no_prior_best else "待继续验证"],
        ["审稿风险", "需要明确区分视频-时间先验、输出校准、生理信号贡献"],
    ]
    add_table(s, 0.75, 1.45, 11.85, 3.25, ["汇报点", "内容"], rows, widths=[2.1, 9.75], font_size=9.1)
    add_box(
        s,
        0.95,
        5.25,
        3.5,
        0.95,
        "不要这样说",
        "“大模型已经强力解码情绪。”",
        RED,
        title_size=12.5,
        body_size=10,
    )
    add_box(
        s,
        4.85,
        5.25,
        6.75,
        0.95,
        "建议这样说",
        "“在强视频时间结构下，我们提出可解释的可信残差融合，用保守信号证据改善跨主体预测。”",
        GREEN,
        title_size=12.5,
        body_size=10,
    )


def append_outline(slide_count: int) -> None:
    additions = [
        "",
        "Added detailed innovation-module slides:",
        f"{slide_count - 4}. CRF-Fusion（可信残差融合）总览图",
        f"{slide_count - 3}. SCRF-BCRF（符号一致贝叶斯可信残差场）结构图",
        f"{slide_count - 2}. CCMI（保守跨模态交集）融合图",
        f"{slide_count - 1}. 数学拼接逻辑",
        f"{slide_count}. 汇报话术：贡献、证据与风险",
    ]
    old = OUTLINE.read_text(encoding="utf-8") if OUTLINE.exists() else ""
    marker = "Added detailed innovation-module slides:"
    if marker in old:
        old = old.split(marker, 1)[0].rstrip() + "\n"
    OUTLINE.write_text(old.rstrip() + "\n" + "\n".join(additions) + "\n", encoding="utf-8")


def main() -> None:
    bcrf = json.loads((RESULT_DIR / "iteration_221_228_bcrf_module_seed2026.json").read_text(encoding="utf-8"))
    summary = json.loads((RESULT_DIR / "three_module_optimization_summary.json").read_text(encoding="utf-8"))
    no_prior_path = RESULT_DIR / "iteration_321_335_no_video_prior_signal.json"
    no_prior = json.loads(no_prior_path.read_text(encoding="utf-8")) if no_prior_path.exists() else None
    best = bcrf["aggregate_results"][0]
    signal_best = summary["ccmi_fusion"][-1]
    no_prior_best = no_prior["aggregate_results"][0] if no_prior else None

    prs = Presentation(PPTX)
    replace_exact_text(prs)
    slide_overall_module(prs)
    slide_scrf_bcrf(prs)
    slide_ccmi(prs)
    slide_math(prs, best, no_prior_best)
    slide_speaker_notes(prs, best, signal_best, no_prior_best)

    for idx, slide in enumerate(prs.slides, 1):
        if idx > len(prs.slides) - 5:
            add_footer(slide, idx)
    prs.save(PPTX)
    append_outline(len(prs.slides))
    print(f"enhanced {PPTX}")
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    main()
