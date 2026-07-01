# -*- coding: utf-8 -*-
"""Label the image2-generated module figure and add it to the MER-PS PPT."""

from __future__ import annotations

import json
import os
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIG_DIR = ROOT / "assets" / "figures"
BASE = FIG_DIR / "crf_fusion_image2_base.png"
LABELED = FIG_DIR / "crf_fusion_image2_labeled.png"
PPTX = Path(os.environ.get("MERPS_PPTX", str(ROOT / "MER_PS_strategy_report.pptx")))
OUTLINE = Path(os.environ.get("MERPS_OUTLINE", str(ROOT / "MER_PS_strategy_report_outline.md")))
RESULT_DIR = ROOT / "experiments" / "results"

FONT = "Microsoft YaHei"
INK = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(248, 250, 252)


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        Path("C:/Windows/Fonts/msyhbd.ttc") if bold else Path("C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
        Path("C:/Windows/Fonts/arial.ttf"),
    ]
    for path in candidates:
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def text_box(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    fill: tuple[int, int, int],
    size: int = 34,
    bold: bool = True,
    anchor: str = "mm",
    align: str = "center",
    stroke: bool = True,
) -> None:
    x, y = xy
    f = font(size, bold)
    kwargs = {}
    if stroke:
        kwargs = {"stroke_width": 4, "stroke_fill": (255, 255, 255)}
    draw.multiline_text(
        (x, y),
        text,
        fill=fill,
        font=f,
        anchor=anchor,
        align=align,
        spacing=5,
        **kwargs,
    )


def make_labeled_figure() -> None:
    img = Image.open(BASE).convert("RGBA")
    draw = ImageDraw.Draw(img)
    w, h = img.size

    # Slide-level title and legend. Coordinates are tuned for the generated 16:9 figure.
    text_box(
        draw,
        (w // 2, 44),
        "CRF-Fusion：可信残差融合框架",
        (15, 23, 42),
        36,
        True,
        stroke=True,
    )
    text_box(
        draw,
        (w // 2, 88),
        "EEG（脑电图）+ fNIRS（功能性近红外光谱） | PatternPrior（视频-时间模式先验） | SCRF-BCRF（可信残差校准）",
        (71, 85, 105),
        21,
        False,
        stroke=True,
    )

    blue = (37, 99, 235)
    teal = (13, 148, 136)
    purple = (124, 58, 237)
    green = (22, 163, 74)
    amber = (217, 119, 6)
    dark = (15, 23, 42)

    # Top physiological signal branch.
    text_box(draw, (230, 214), "EEG / fNIRS\n生理信号", blue, 28)
    text_box(draw, (615, 214), "信号特征\nφ_s(signal)", blue, 27)
    text_box(draw, (930, 214), "低维压缩\nPCA / PLS", blue, 27)
    text_box(draw, (1245, 214), "信号专家\nΔ_sig", blue, 27)

    # Middle prior branch.
    text_box(draw, (220, 472), "sample_id\n视频 / 时间", teal, 28)
    text_box(draw, (610, 472), "PatternPrior\np0(v,t)", teal, 27)
    text_box(draw, (930, 472), "状态轨迹\nslope(p0)", teal, 27)
    text_box(draw, (1240, 472), "CCMI\n保守交集", teal, 27)

    # Bottom credible residual branch.
    text_box(draw, (225, 724), "OOF residual\n折外残差 r", purple, 28)
    text_box(draw, (610, 724), "SCRF\n符号一致", purple, 27)
    text_box(draw, (930, 724), "BCRF\n可信权重", purple, 27)
    text_box(draw, (1245, 724), "Δ_cal\n校准残差", purple, 27)

    # Fusion and output.
    text_box(draw, (1442, 507), "融合门控\nλ·g_ccmi", amber, 25)
    text_box(draw, (1595, 500), "输出\nV / A", green, 24)
    text_box(
        draw,
        (1290, 614),
        "ŷ = clip(p0 + Δ_cal + λ·g_ccmi·Δ_sig, 1, 255)",
        dark,
        25,
        True,
        stroke=True,
    )

    # Branch labels at the left margin.
    text_box(draw, (52, 145), "生理证据支路", blue, 23, True, anchor="lm")
    text_box(draw, (52, 394), "视频先导支路", teal, 23, True, anchor="lm")
    text_box(draw, (52, 650), "可信校准支路", purple, 23, True, anchor="lm")

    FIG_DIR.mkdir(parents=True, exist_ok=True)
    img.convert("RGB").save(LABELED, quality=96)


def set_font(run, size=None, bold=False, color=INK):
    run.font.name = FONT
    run.font.bold = bold
    run.font.color.rgb = color
    if size is not None:
        run.font.size = Pt(size)


def add_bg(slide):
    rect = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG
    rect.line.fill.background()
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.1), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = title
    set_font(r, 24, True, INK)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.96), Inches(12.0), Inches(0.32))
        stf = sub.text_frame
        stf.clear()
        r2 = stf.paragraphs[0].add_run()
        r2.text = subtitle
        set_font(r2, 10.5, False, MUTED)


def add_footer(slide, idx: int):
    box = slide.shapes.add_textbox(Inches(11.8), Inches(7.06), Inches(0.9), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{idx:02d}"
    set_font(r, 8, False, MUTED)


def add_text(slide, x, y, w, h, text, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, bold, color)


def add_table(slide, x, y, w, h, headers, rows, widths=None, font_size=8.7):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = ""
        r = p.add_run()
        r.text = header
        set_font(r, font_size, True, WHITE)
    for i, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            cell = table.cell(i, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(241, 245, 249)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            p.text = ""
            r = p.add_run()
            r.text = str(value)
            set_font(r, font_size, False, INK)


def new_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle)
    return slide


def load_results():
    best = json.loads((RESULT_DIR / "iteration_221_228_bcrf_module_seed2026.json").read_text(encoding="utf-8"))[
        "aggregate_results"
    ][0]
    no_prior = json.loads((RESULT_DIR / "iteration_321_335_no_video_prior_signal.json").read_text(encoding="utf-8"))[
        "aggregate_results"
    ][0]
    official = json.loads((RESULT_DIR / "official_demo_eval.json").read_text(encoding="utf-8"))
    signal_best = json.loads((RESULT_DIR / "three_module_optimization_summary.json").read_text(encoding="utf-8"))[
        "ccmi_fusion"
    ][-1]
    return best, no_prior, official, signal_best


def add_picture_slide(prs):
    slide = new_slide(
        prs,
        "Image2 生成模块图：CRF-Fusion（可信残差融合）",
        "该图由 image2 生成科研风格底图，再叠加准确术语、公式和中文解释",
    )
    slide.shapes.add_picture(str(LABELED), Inches(0.45), Inches(1.28), width=Inches(12.45))
    add_text(
        slide,
        0.78,
        6.82,
        11.9,
        0.25,
        "图中三条支路分别对应：直接生理证据、视频/时间先导、可信残差校准；两种提交路线会分别启用不同支路组合。",
        9.2,
        False,
        MUTED,
    )


def add_submission_slides(prs):
    best, no_prior, official, signal_best = load_results()

    s = new_slide(
        prs,
        "提交路线 A：视频/时间先导高分模型",
        "目标是 leaderboard 分数；允许利用 sample_id 中的 video/time 结构，但必须说明其先验属性",
    )
    add_table(
        s,
        0.7,
        1.42,
        11.95,
        2.65,
        ["模块", "启用内容", "作用", "当前证据"],
        [
            ["PatternPrior（视频-时间模式先验）", "p0(video,time)", "学习每个视频每秒的平均情绪轨迹", "PatternPrior_098: 28.7462"],
            ["SCRF-BCRF（可信残差校准）", "Δ_cal", "只修正可信系统性残差，主要修 valence", f"{best['method']}: {best['overall_mae']:.4f}"],
            ["CCMI（保守跨模态交集）", "λ·g·Δ_sig", "作为小幅生理残差候选", f"Best signal path: {signal_best['overall_mae']:.4f}"],
        ],
        widths=[3.05, 2.2, 4.1, 2.6],
        font_size=8.5,
    )
    add_text(
        s,
        0.85,
        4.65,
        11.3,
        0.95,
        "提交形式：输出 raw [1,255] 的 valence（效价）和 arousal（唤醒度）回归值。当前推荐主提交使用 222_BCRF_onSCRF；如果后续完成 per-sample cache，再尝试叠加小权重 CCMI。",
        12,
        False,
        INK,
    )
    add_text(
        s,
        0.85,
        5.85,
        11.3,
        0.72,
        "汇报时要强调：这条路线是比赛高分路线，利用了公开训练集中可复现的视频/时间轨迹结构，不等价于纯生理信号解码。",
        11.5,
        True,
        RED,
    )

    s = new_slide(
        prs,
        "提交路线 B：直接生理信号模型",
        "目标是验证 EEG+fNIRS 本身的可解码性；不使用 VideoTimeMean / PatternPrior / video-time cell",
    )
    add_table(
        s,
        0.7,
        1.42,
        11.95,
        2.75,
        ["方法", "Overall MAE", "Valence MAE", "Arousal MAE", "解释"],
        [
            ["Official ASAC demo（官方演示）", f"{official['overall_mae']:.4f}", f"{official['valence_mae']:.4f}", f"{official['arousal_mae']:.4f}", "官方式信号模型参考"],
            ["Center128 no-prior", "47.5663", "52.1980", "42.9346", "无信号、无先验"],
            [no_prior["method"], f"{no_prior['overall_mae']:.4f}", f"{no_prior['valence_mae']:.4f}", f"{no_prior['arousal_mae']:.4f}", "当前最佳直接生理路线"],
        ],
        widths=[4.05, 1.35, 1.35, 1.35, 3.85],
        font_size=8.4,
    )
    add_text(
        s,
        0.85,
        4.7,
        11.35,
        0.95,
        "提交形式仍是回归预测，不改成离散分类；这里的“直接生理”指模型只读 EEG/fNIRS 特征，不读训练标签形成的视频-时间先验。",
        12,
        False,
        INK,
    )
    add_text(
        s,
        0.85,
        5.9,
        11.35,
        0.72,
        f"当前结论：直接生理模型比 Center128 改善约 {47.5663 - no_prior['overall_mae']:.4f} MAE，但比高分先导路线差约 {no_prior['overall_mae'] - best['overall_mae']:.4f} MAE；它适合作为论文中的“去先验生理证据”对照提交。",
        11.5,
        True,
        TEAL,
    )

    s = new_slide(
        prs,
        "两种提交路线如何同时使用",
        "一个追求比赛分数，一个支撑论文可信性和消融分析",
    )
    add_table(
        s,
        0.7,
        1.42,
        11.95,
        3.55,
        ["维度", "路线 A：视频/时间先导", "路线 B：直接生理"],
        [
            ["输入", "sample_id + EEG/fNIRS，可使用 video/time 统计先验", "只用 EEG/fNIRS 特征，不用 video/time label prior"],
            ["目标", "争取 leaderboard 最低 MAE", "验证生理信号本身是否可解码"],
            ["当前最好", f"{best['overall_mae']:.4f}", f"{no_prior['overall_mae']:.4f}"],
            ["论文作用", "证明可信残差校准能提升强基线", "回应审稿人对先验泄漏/刺激记忆的质疑"],
            ["风险", "容易被认为依赖视频刺激模式", "分数较弱，需要解释为去先验对照"],
        ],
        widths=[1.8, 5.05, 5.1],
        font_size=8.8,
    )
    add_text(
        s,
        0.85,
        5.55,
        11.35,
        0.78,
        "建议提交策略：主榜提交路线 A；同时保留路线 B 作为备选/消融提交，并在论文或答辩中明确报告两者差距。",
        12.2,
        True,
        GREEN,
    )


def update_outline(slide_count: int):
    marker = "Image2 generated module figure and two-route submission slides:"
    old = OUTLINE.read_text(encoding="utf-8") if OUTLINE.exists() else ""
    if marker in old:
        old = old.split(marker, 1)[0].rstrip() + "\n"
    lines = [
        "",
        marker,
        f"{slide_count - 3}. Image2 模块图：CRF-Fusion（可信残差融合）",
        f"{slide_count - 2}. 提交路线 A：视频/时间先导高分模型",
        f"{slide_count - 1}. 提交路线 B：直接生理信号模型",
        f"{slide_count}. 两种提交路线如何同时使用",
    ]
    OUTLINE.write_text(old.rstrip() + "\n" + "\n".join(lines) + "\n", encoding="utf-8")


def main():
    make_labeled_figure()
    prs = Presentation(PPTX)
    add_picture_slide(prs)
    add_submission_slides(prs)
    for idx, slide in enumerate(prs.slides, 1):
        if idx > len(prs.slides) - 4:
            add_footer(slide, idx)
    prs.save(PPTX)
    update_outline(len(prs.slides))
    print(f"wrote {LABELED}")
    print(f"updated {PPTX}")
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    main()
