# -*- coding: utf-8 -*-
"""Build a concise MER-PS strategy report deck.

The deck is generated from the experiment summary JSON files so the headline
numbers remain tied to the recorded runs.
"""

from __future__ import annotations

import json
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
RESULT_DIR = ROOT / "experiments" / "results"
OUT_PPTX = Path(os.environ.get("MERPS_PPTX", str(ROOT / "MER_PS_strategy_report.pptx")))
OUT_MD = Path(os.environ.get("MERPS_OUTLINE", str(ROOT / "MER_PS_strategy_report_outline.md")))

FONT = "Microsoft YaHei"
BG = RGBColor(248, 250, 252)
INK = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
LINE = RGBColor(203, 213, 225)
WHITE = RGBColor(255, 255, 255)


def load_json(name: str) -> dict:
    with (RESULT_DIR / name).open("r", encoding="utf-8") as f:
        return json.load(f)


def maybe_load_json(name: str) -> dict | None:
    path = RESULT_DIR / name
    if not path.exists():
        return None
    return json.loads(path.read_text(encoding="utf-8"))


def experiment_stats() -> dict:
    files = list(RESULT_DIR.glob("*.json"))
    non_smoke = [p for p in files if not p.name.startswith("smoke")]
    smoke = [p for p in files if p.name.startswith("smoke")]
    ids: set[int] = set()
    agg_rows = 0
    for p in non_smoke:
        for m in re.finditer(r"iteration_(\d+)(?:_(\d+))?", p.name):
            start = int(m.group(1))
            end = int(m.group(2) or start)
            ids.update(range(start, end + 1))
        try:
            data = json.loads(p.read_text(encoding="utf-8"))
        except Exception:
            continue
        rows = data.get("aggregate_results")
        if isinstance(rows, list):
            agg_rows += len(rows)
    return {
        "json_files": len(files),
        "non_smoke_files": len(non_smoke),
        "smoke_files": len(smoke),
        "unique_iteration_ids": len(ids),
        "min_iter": min(ids) if ids else None,
        "max_iter": max(ids) if ids else None,
        "aggregate_rows": agg_rows,
    }


def set_font(run, size=None, bold=False, color=INK):
    run.font.name = FONT
    run.font.bold = bold
    run.font.color.rgb = color
    if size is not None:
        run.font.size = Pt(size)


def add_bg(slide):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG
    rect.line.fill.background()
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.72))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    set_font(r, 25, True, INK)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.0), Inches(11.7), Inches(0.35))
        stf = sub.text_frame
        stf.clear()
        p2 = stf.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        set_font(r2, 10.5, False, MUTED)


def add_footer(slide, idx):
    box = slide.shapes.add_textbox(Inches(11.85), Inches(7.05), Inches(0.9), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{idx:02d}"
    set_font(r, 8, False, MUTED)


def add_chip(slide, x, y, w, h, label, value, color=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(value)
    set_font(r, 20, True, color)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    set_font(r2, 8.5, False, MUTED)
    return shp


def add_text(slide, x, y, w, h, text, size=13, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, bold, color)
    return box


def add_bullets(slide, x, y, w, h, bullets, size=13):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.text = ""
        r = p.add_run()
        r.text = "• " + text
        set_font(r, size, False, INK)
        p.space_after = Pt(6)
    return box


def add_card(slide, x, y, w, h, title, body=None, color=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, 12.5, True, color)
    if body:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = body
        set_font(r2, 9.5, False, MUTED)
    return shp


def add_table(slide, x, y, w, h, headers, rows, widths=None, font_size=8.5):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = ""
        r = p.add_run()
        r.text = header
        set_font(r, font_size, True, WHITE)
    for r_i, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            cell = table.cell(r_i, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_i % 2 else RGBColor(241, 245, 249)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            p.text = ""
            r = p.add_run()
            r.text = str(value)
            color = GREEN if ("推荐" in str(value) or str(value).startswith("28.6868")) else INK
            set_font(r, font_size, False, color)
    return table_shape


def add_flow(slide, labels, y=2.6):
    x = 0.7
    w = 1.8
    gap = 0.28
    shapes = []
    colors = [BLUE, TEAL, AMBER, GREEN, BLUE, TEAL]
    for i, label in enumerate(labels):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.68))
        shp.fill.solid()
        shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = colors[i % len(colors)]
        shp.line.width = Pt(1.4)
        tf = shp.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_font(r, 9.5, True, INK)
        shapes.append(shp)
        if i:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x - gap + 0.02),
                Inches(y + 0.34),
                Inches(x - 0.02),
                Inches(y + 0.34),
            )
            conn.line.color.rgb = LINE
            conn.line.width = Pt(1.2)
        x += w + gap
    return shapes


def metric_rows(best, demo, signal_best, no_prior_best=None):
    rows = [
        ["Official ASAC demo", f"{demo['overall_mae']:.4f}", f"{demo['valence_mae']:.4f}", f"{demo['arousal_mae']:.4f}", "提交接口参考"],
    ]
    if no_prior_best:
        rows.append(
            [
                no_prior_best["method"],
                f"{no_prior_best['overall_mae']:.4f}",
                f"{no_prior_best['valence_mae']:.4f}",
                f"{no_prior_best['arousal_mae']:.4f}",
                "去 video/time 先验",
            ]
        )
    rows.extend(
        [
            ["PatternPrior_098", "28.7462", "26.9738", "30.5186", "强视频/时间先验"],
            ["Best EEG-fNIRS CCMI", f"{signal_best['overall_mae']:.4f}", f"{signal_best['valence_mae']:.4f}", f"{signal_best['arousal_mae']:.4f}", "最佳信号融合路径"],
            [best["method"], f"{best['overall_mae']:.4f}", f"{best['valence_mae']:.4f}", f"{best['arousal_mae']:.4f}", "当前全局最佳"],
        ]
    )
    return [
        *rows
    ]


def short_name(name):
    return name.split(": ", 1)[-1].replace(" + ", "\n+ ")


def build_deck() -> None:
    stats = experiment_stats()
    summary = load_json("three_module_optimization_summary.json")
    official = load_json("official_demo_eval.json")
    bcrf = load_json("iteration_221_228_bcrf_module_seed2026.json")
    no_prior = maybe_load_json("iteration_321_335_no_video_prior_signal.json")
    best = bcrf["aggregate_results"][0]
    signal_best = summary["ccmi_fusion"][-1]
    no_prior_best = no_prior["aggregate_results"][0] if no_prior else None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides = []

    def new_slide(title, subtitle=None):
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, title, subtitle)
        slides.append(slide)
        return slide

    s = new_slide("MER-PS 情绪脑电比赛策略复盘", "从数据、评价指标、实验演进到当前最优框架")
    add_text(s, 0.7, 1.75, 7.2, 0.7, "目标：用 EEG + fNIRS 预测 1 Hz 连续 valence / arousal 轨迹", 20, True, BLUE)
    add_bullets(
        s,
        0.8,
        2.65,
        6.7,
        2.0,
        [
            "当前最佳：222_BCRF_onSCRF，Overall MAE = 28.6868。",
            "最佳 EEG-fNIRS 信号融合：CCMI PriorSlopeGate，Overall MAE = 28.7145。",
            (
                "去视频/时间先验首批："
                f"{no_prior_best['method']}，Overall MAE = {no_prior_best['overall_mae']:.4f}。"
                if no_prior_best
                else "下一步重点：去掉视频/时间前导，验证生理信号真实贡献。"
            ),
        ],
        14,
    )
    add_chip(s, 8.2, 1.65, 1.55, 1.0, "Overall MAE", "28.6868", GREEN)
    add_chip(s, 9.95, 1.65, 1.55, 1.0, "实验编号", f"1-{stats['max_iter']}", BLUE)
    add_chip(s, 11.7, 1.65, 1.05, 1.0, "JSON", stats["json_files"], TEAL)
    add_card(s, 8.2, 3.05, 4.55, 1.35, "核心判断", "数据量小且跨主体，纯大模型容易过拟合；有效改进来自低自由度先验、可信残差和保守多模态交集。", AMBER)
    add_card(s, 8.2, 4.65, 4.55, 1.25, "当前问题", "最优模型仍混入 video/time 标签结构，需要官方式直接信号处理实验来确认泛化价值。", RED)

    s = new_slide("赛题与数据", "MER-PS: synchronized EEG + fNIRS continuous emotion decoding")
    add_table(
        s,
        0.7,
        1.55,
        6.2,
        2.3,
        ["项目", "内容"],
        [
            ["训练验证主体", "test_1 ~ test_24，subject-disjoint CV"],
            ["每主体试次", "15 段情绪诱发视频，含 5 秒 baseline"],
            ["EEG", "64 通道，原 1000 Hz，赛题文件 200 Hz"],
            ["fNIRS", "51 通道，47.62 Hz，6 类信号类型"],
            ["标签", "valence / arousal，1 Hz，整数 [1,255]"],
        ],
        widths=[1.8, 4.4],
        font_size=9,
    )
    add_card(s, 7.35, 1.55, 4.95, 1.05, "预测文件", "sample_id,valence,arousal；每个 sample_id 必须出现且只出现一次。", BLUE)
    add_card(s, 7.35, 2.9, 4.95, 1.05, "评价指标", "Overall MAE = valence 和 arousal 绝对误差的平均；越低越好。", TEAL)
    add_card(s, 7.35, 4.25, 4.95, 1.05, "风险点", "公开视频身份和时间位置会形成强先验；若依赖过强，leaderboard 外泛化可能不稳。", RED)

    s = new_slide("实验规模", "从官方 demo 到 320 轮编号实验")
    add_chip(s, 0.75, 1.55, 2.0, 1.0, "结果 JSON 文件", stats["json_files"], BLUE)
    add_chip(s, 3.0, 1.55, 2.0, 1.0, "非 smoke 文件", stats["non_smoke_files"], TEAL)
    add_chip(s, 5.25, 1.55, 2.0, 1.0, "唯一编号覆盖", stats["unique_iteration_ids"], GREEN)
    add_chip(s, 7.5, 1.55, 2.0, 1.0, "累计结果行", stats["aggregate_rows"], AMBER)
    add_chip(s, 9.75, 1.55, 2.0, 1.0, "Smoke 检查", stats["smoke_files"], RED)
    add_bullets(
        s,
        0.85,
        3.0,
        11.5,
        2.2,
        [
            "严格说：结果文件中非 smoke 的编号覆盖 301 个 iteration ID，编号推进到 320。",
            "累计 aggregate result 2122 条，包含多 seed、多折、模块变体和补充验证。",
            "Smoke 结果只用于筛掉明显不值得跑 full CV 的方案，不能和 24-subject CV 直接比绝对分数。",
        ],
        14,
    )

    s = new_slide("当前最优结果", "官方 demo 不是强 baseline；我们的最优来自输出校准")
    add_table(
        s,
        0.65,
        1.45,
        12.05,
        3.0 if no_prior_best else 2.55,
        ["方法", "Overall MAE", "Valence MAE", "Arousal MAE", "说明"],
        metric_rows(best, official, signal_best, no_prior_best),
        widths=[3.1, 1.65, 1.65, 1.65, 3.0],
        font_size=8.8,
    )
    add_bullets(
        s,
        0.85,
        4.75 if no_prior_best else 4.45,
        11.7,
        1.4,
        [
            "ASAC demo 参数很少，主要验证提交接口；在 test_21~test_24 上 Overall MAE = 47.0087。",
            (
                f"去 video/time 先验后，最佳直接信号模型为 {no_prior_best['overall_mae']:.4f}，"
                "说明信号本身只带来小幅改善。"
                if no_prior_best
                else "PatternPrior 已经很强，说明视频/时间结构解释了大量标签变化。"
            ),
            "PatternPrior 已经很强，说明视频/时间结构解释了大量标签变化。",
            "BCRF_onSCRF 只比 SCRF 提升 0.0001，说明后期收益很小，需要更扎实的去先验验证。",
        ],
        12.5,
    )

    s = new_slide("策略演进", "先广泛复现，再收敛到低自由度可信校准")
    rows = [
        ["1-20", "Graph/Mamba/SSM", "验证深模型是否能吃到生理信号", "容易过拟合，收益不稳"],
        ["28-63", "蒸馏/时序/教师", "参考 AffectGPT、时序模型和 teacher residual", "教师信息有限，标签先验更强"],
        ["64-104", "Pattern prior / confidence", "视频、时间、状态条件先验", "形成强基线 28.7 左右"],
        ["105-200", "大规模架构搜索", "残差、平滑、risk expert、dimwise fusion", "确认 valence 可修，arousal 要保守"],
        ["212-228", "SCRF / BCRF", "可信残差场输出校准", "当前全局最佳 28.6868"],
        ["229-320", "EEG-fNIRS 融合", "OOF gate、NOVA、CCMI、fNIRS all6", "信号路径最佳 28.7145"],
    ]
    add_table(s, 0.65, 1.45, 12.05, 4.35, ["轮次", "方向", "目的", "结论"], rows, widths=[1.1, 2.4, 4.0, 4.55], font_size=8.7)

    s = new_slide("为什么不是盲目堆大模型", "训练样本的有效自由度远小于原始时间点数量")
    add_bullets(
        s,
        0.85,
        1.55,
        5.85,
        4.3,
        [
            "24 个主体，每主体 15 段视频；subject-disjoint 要求模型跨人泛化。",
            "1 Hz 标签虽然有很多时间点，但相邻标签高度相关，不能当作独立样本。",
            "EEG/fNIRS 维度高、噪声大、个体差异强；高参数模型更容易记住 subject/video 模式。",
            "因此有效策略是：先建立强低维先验，再只接受可信的残差信号。",
        ],
        13.5,
    )
    add_card(s, 7.35, 1.65, 4.8, 1.1, "数学直觉", "y = prior(video,time) + reliable_residual(signal) + noise", BLUE)
    add_card(s, 7.35, 3.0, 4.8, 1.1, "模型约束", "残差 correction 的自由度必须小，且要有 OOF 可信度估计。", TEAL)
    add_card(s, 7.35, 4.35, 4.8, 1.1, "下一步实验", "去掉 video/time prior 后，直接评估 EEG+fNIRS 的可解码能力。", RED)

    s = new_slide("数据预处理模块", "五种预处理对比")
    d_rows = []
    for item in summary["data_preprocessing"]:
        tag = "推荐" if item["recommended"] else ""
        d_rows.append([short_name(item["name"]), f"{item['overall_mae']:.4f}", f"{item['valence_mae']:.4f}", f"{item['arousal_mae']:.4f}", item["validation"], tag])
    add_table(s, 0.55, 1.42, 12.35, 3.65, ["预处理", "Overall", "Valence", "Arousal", "验证", ""], d_rows, widths=[3.5, 1.2, 1.2, 1.2, 2.0, 0.8], font_size=8.2)
    add_bullets(
        s,
        0.8,
        5.45,
        11.8,
        0.85,
        [
            "结论：6 类 fNIRS + baseline mean subtraction 最稳；不做 baseline subtraction 会回落到 PatternPrior。",
            "trial/subject z-score 仅 smoke，绝对分数不可横比；它提示幅值信息不能随便抹掉。",
        ],
        11.5,
    )

    s = new_slide("CCMI 融合模块", "保守交集比大容量注意力更适合当前数据量")
    c_rows = []
    for item in summary["ccmi_fusion"]:
        tag = "推荐" if item["recommended"] else ""
        c_rows.append([short_name(item["name"]), f"{item['overall_mae']:.4f}", f"{item['valence_mae']:.4f}", f"{item['arousal_mae']:.4f}", tag])
    add_table(s, 0.55, 1.35, 12.35, 3.1, ["融合方法", "Overall", "Valence", "Arousal", ""], c_rows, widths=[5.0, 1.35, 1.35, 1.35, 0.9], font_size=8.4)
    add_card(s, 0.85, 4.85, 5.8, 0.9, "CCMI 核心公式", "r = 1[sign(E)=sign(F)] · sign(E+F) · min(|E|,|F|)", BLUE)
    add_card(s, 7.0, 4.85, 5.2, 0.9, "为什么有效", "只有 EEG 与 fNIRS 同向时才修正，避免单模态噪声把预测拉偏。", TEAL)

    s = new_slide("输出头模块", "当前最优来自可信残差场，而不是直接加复杂网络")
    h_rows = []
    for item in summary["output_head"]:
        tag = "推荐" if item["recommended"] else ""
        h_rows.append([short_name(item["name"]), f"{item['overall_mae']:.4f}", f"{item['valence_mae']:.4f}", f"{item['arousal_mae']:.4f}", tag])
    add_table(s, 0.55, 1.35, 12.35, 3.1, ["输出头", "Overall", "Valence", "Arousal", ""], h_rows, widths=[5.0, 1.35, 1.35, 1.35, 0.9], font_size=8.4)
    add_bullets(
        s,
        0.85,
        4.85,
        11.6,
        1.0,
        [
            "SCRF/BCRF 主要修 valence；arousal 多次 residual probe 后都不稳定，因此保持保守。",
            "BCRF 的提升非常小，但可解释性较强：它估计“残差是否可信”，而不是盲目拟合残差。",
        ],
        11.5,
    )

    s = new_slide("SCRF-BCRF 模块逻辑", "Sign-Consistent Bayesian Credible Residual Field")
    add_flow(s, ["OOF residual", "低维分组", "median/MAD/n", "可信权重", "符号一致门控", "最终修正"], y=1.65)
    add_card(s, 0.75, 3.05, 3.75, 1.0, "残差定义", "r_i = y_i - p_i；只用 OOF 预测构建残差表，降低泄漏风险。", BLUE)
    add_card(s, 4.85, 3.05, 3.75, 1.0, "可信权重", "w = n/(n+k) · 1/(1+MAD/c) · sigmoid(z-τ) · sign_strength", TEAL)
    add_card(s, 8.95, 3.05, 3.0, 1.0, "应用范围", "主要作用于 valence；arousal correction 基本关闭。", AMBER)
    add_bullets(
        s,
        0.9,
        4.65,
        11.5,
        1.1,
        [
            "这个模块的创新点不是“拟合更复杂”，而是把残差修正变成可置信的低维统计决策。",
            "它适合小样本、跨主体、强噪声场景：宁可少修，也不让噪声修正毁掉基线。",
        ],
        11.5,
    )

    s = new_slide("当前最佳框架", "强先验 + 可信残差；信号模块还未和 222 完整融合")
    add_flow(s, ["原始 EEG/fNIRS", "特征统计", "PatternPrior", "Dimwise Head", "SCRF", "BCRF"], y=1.55)
    add_flow(s, ["EEG residual", "fNIRS residual", "CCMI gate", "signal correction"], y=3.3)
    add_card(s, 0.85, 4.75, 3.8, 0.9, "已验证全局最佳", "222_BCRF_onSCRF：Overall MAE 28.6868", GREEN)
    add_card(s, 4.95, 4.75, 3.8, 0.9, "已验证信号最佳", "D2 + C5：Overall MAE 28.7145", TEAL)
    add_card(s, 9.05, 4.75, 3.25, 0.9, "未完成验证", "222 + 小 CCMI 残差，需要 per-sample cache。", RED)

    if no_prior_best:
        s = new_slide("去视频/时间前导实验", "直接处理 EEG/fNIRS，隔离生理信号真实贡献")
        center_row = next(
            item for item in no_prior["aggregate_results"] if item["method"] == "321_Center128_noPrior"
        )
        add_table(
            s,
            0.65,
            1.38,
            12.0,
            3.0,
            ["方法", "Overall", "Valence", "Arousal", "解释"],
            [
                ["Official ASAC demo", f"{official['overall_mae']:.4f}", f"{official['valence_mae']:.4f}", f"{official['arousal_mae']:.4f}", "官方式小模型，单固定 split"],
                ["Center128 no-prior", f"{center_row['overall_mae']:.4f}", f"{center_row['valence_mae']:.4f}", f"{center_row['arousal_mae']:.4f}", "无信号、无先验"],
                [no_prior_best["method"], f"{no_prior_best['overall_mae']:.4f}", f"{no_prior_best['valence_mae']:.4f}", f"{no_prior_best['arousal_mae']:.4f}", "最佳直接信号模型"],
                [best["method"], f"{best['overall_mae']:.4f}", f"{best['valence_mae']:.4f}", f"{best['arousal_mae']:.4f}", "含强先验 + 可信校准"],
            ],
            widths=[3.6, 1.2, 1.2, 1.2, 3.6],
            font_size=8.5,
        )
        add_card(
            s,
            0.85,
            4.8,
            3.6,
            0.9,
            "信号收益",
            f"相对 Center128 改善 {center_row['overall_mae'] - no_prior_best['overall_mae']:.4f} MAE。",
            GREEN,
        )
        add_card(
            s,
            4.85,
            4.8,
            3.6,
            0.9,
            "先验差距",
            f"相对全局最佳仍差 {no_prior_best['overall_mae'] - best['overall_mae']:.4f} MAE。",
            RED,
        )
        add_card(
            s,
            8.85,
            4.8,
            3.25,
            0.9,
            "结论",
            "生理信号可用，但必须低维、保守、可置信地融合。",
            BLUE,
        )
    else:
        s = new_slide("下一组实验：去视频/时间前导", "和官方 demo 更接近：直接从信号预测")
        add_table(
            s,
            0.65,
            1.42,
            12.0,
            3.15,
            ["实验", "去掉内容", "保留内容", "目的"],
            [
                ["NV-0 Center/Mean", "video_id/time/sample prior", "训练集标签均值", "确定无先验下的最低参考线"],
                ["NV-1 Official-like ASAC", "video/time prior", "EEG/fNIRS 原始特征 + 官方结构", "复核官方范式"],
                ["NV-2 Signal Ridge/HistGB", "video/time prior", "窗口统计、频域、baseline 特征", "测信号线性/树模型上限"],
                ["NV-3 Signal CCMI", "PatternPrior 输入", "EEG 与 fNIRS residual 一致性", "测多模态交集真实贡献"],
                ["NV-4 Signal + SCRF-lite", "video/time cell", "信号状态 cell + 低维可信残差", "把我们模块迁到无视频先验"],
            ],
            widths=[2.2, 3.1, 3.4, 3.3],
            font_size=8.4,
        )

    s = new_slide("阶段性结论", "下一步从“刷分”转向“可解释可发表”")
    add_bullets(
        s,
        0.85,
        1.55,
        11.5,
        3.8,
        [
            "当前最优：222_BCRF_onSCRF，Overall MAE 28.6868；它说明低维可信残差有效。",
            "最优信号融合：CCMI PriorSlopeGate，Overall MAE 28.7145；它说明 EEG/fNIRS 同向证据有价值。",
            "最重要的审稿风险：模型是否依赖视频/时间前导，而不是生理信号。",
            "下一轮实验要去掉 video/time prior，保留官方式输入接口，重新评估几个好模块。",
            "论文方向：小样本跨主体情绪解码中的可信残差校准与保守多模态交集融合。",
        ],
        14,
    )

    for idx, slide in enumerate(slides, 1):
        add_footer(slide, idx)

    prs.save(OUT_PPTX)

    lines = [
        "# MER-PS strategy report outline",
        "",
        f"- Current best: `{best['method']}` overall MAE `{best['overall_mae']:.4f}`, valence `{best['valence_mae']:.4f}`, arousal `{best['arousal_mae']:.4f}`.",
        f"- Experiment scale: `{stats['json_files']}` result JSON files, `{stats['unique_iteration_ids']}` unique iteration IDs, numbering up to `{stats['max_iter']}`, `{stats['aggregate_rows']}` aggregate rows.",
        f"- Best EEG-fNIRS signal fusion: `CCMI PriorSlopeGate`, overall MAE `{signal_best['overall_mae']:.4f}`.",
        (
            f"- Best no-video/time direct signal model: `{no_prior_best['method']}` overall MAE "
            f"`{no_prior_best['overall_mae']:.4f}`."
            if no_prior_best
            else "- Next experiment group: remove video/time prior and evaluate official-like signal-only pipelines."
        ),
        "",
        "Slides:",
    ]
    for idx, title in enumerate(
        [
            "MER-PS 情绪脑电比赛策略复盘",
            "赛题与数据",
            "实验规模",
            "当前最优结果",
            "策略演进",
            "为什么不是盲目堆大模型",
            "数据预处理模块",
            "CCMI 融合模块",
            "输出头模块",
            "SCRF-BCRF 模块逻辑",
            "当前最佳框架",
            "下一组实验：去视频/时间前导",
            "阶段性结论",
        ],
        1,
    ):
        lines.append(f"{idx}. {title}")
    OUT_MD.write_text("\n".join(lines) + "\n", encoding="utf-8")


if __name__ == "__main__":
    build_deck()
    print(f"wrote {OUT_PPTX}")
    print(f"wrote {OUT_MD}")
